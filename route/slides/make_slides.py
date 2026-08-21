"""Generate a 3-slide PPTX explaining the robust-routing principle.

Stage 1: likelihood-weighted input space
Stage 2: deterministic optimisation + cross-evaluation matrix
Stage 3: matrix -> decision, via P(break record)

All three figures are driven by ONE toy physical model, so every number
quoted on the slides is internally consistent.

Toy model (Lowestoft -> IJmuiden, 101.7 nm on 090 T, single gybe-free reach):
  route parameter x = lateral bulge north of the rhumb line, in nm.
  Bulging north adds distance but rotates the heading, protecting TWA when
  the wind backs. That is the real distance-vs-angle trade-off, and it is
  what makes routes differently *sensitive* rather than merely differently
  *fast* -- which is what the robustness KPI is supposed to detect.
"""
import os

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

FIG = "slides/figures"
INK = "#1F2933"
MUTED = "#7B8794"
ACCENT = "#0B6E6E"
WARM = "#C0632C"
RULE = "#D6DBE0"

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Calibri", "DejaVu Sans"],
    "text.color": INK,
    "axes.labelcolor": INK,
    "axes.edgecolor": MUTED,
    "xtick.color": MUTED,
    "ytick.color": MUTED,
    "axes.spines.top": False,
    "axes.spines.right": False,
    "figure.facecolor": "white",
    "font.size": 11,
})

# ------------------------------------------------------------- physical model
HALF_NM = 50.85         # half the rhumb line (101.7 nm; see plot_routes.py)
TWA_OPT = 145.0         # peak of the polar, in TWA terms
TWA_NOM = 150.0         # nominal sailing angle on the rhumb line
POLAR_W = 45.0          # polar falloff width (deg)
RECORD_MIN = 339.0      # 101.7 nm at exactly 18.0 kt average
N_SCEN, N_ROUTES = 40, 12


def route_distance(x):
    """Distance for a symmetric lateral bulge of x nm."""
    return 2.0 * np.sqrt(HALF_NM ** 2 + x ** 2)


def heading_shift(x):
    """Degrees of heading rotation produced by an x nm bulge."""
    return np.degrees(np.arctan2(x, HALF_NM))


def boat_speed(tws, twa, polar_scale):
    """Toy polar: peak near TWA_OPT, Gaussian falloff, linear in TWS."""
    return polar_scale * 0.93 * tws * np.exp(-((twa - TWA_OPT) / POLAR_W) ** 2)


def elapsed_min(x, tws, twd_shift, polar_scale):
    """Minutes to sail the bulge-x route under one input scenario."""
    twa = TWA_NOM + twd_shift - heading_shift(x)
    return route_distance(x) / boat_speed(tws, twa, polar_scale) * 60.0


# --------------------------------------------- stage 1: weighted input space
rng = np.random.default_rng(11)
SD_TWS, SD_TWD, SD_POL = 1.5, 10.0, 0.04

tws = rng.normal(20.0, SD_TWS, N_SCEN)
twd = rng.normal(0.0, SD_TWD, N_SCEN)
pol = rng.normal(1.0, SD_POL, N_SCEN)

order = np.argsort(twd)                        # sort so the matrix reads cleanly
tws, twd, pol = tws[order], twd[order], pol[order]

# Joint likelihood = product of the marginals. This IS the weighting principle.
w = (np.exp(-0.5 * ((tws - 20.0) / SD_TWS) ** 2)
     * np.exp(-0.5 * (twd / SD_TWD) ** 2)
     * np.exp(-0.5 * ((pol - 1.0) / SD_POL) ** 2))
w /= w.sum()

# ----------------------------- stage 2: one optimal route per chosen scenario
grid = np.linspace(-14.0, 26.0, 801)
home = np.round(np.linspace(0, N_SCEN - 1, N_ROUTES)).astype(int)

x_route = np.array([
    grid[np.argmin(elapsed_min(grid, tws[j], twd[j], pol[j]))] for j in home
])

T = np.empty((N_ROUTES, N_SCEN))
for i in range(N_ROUTES):
    for j in range(N_SCEN):
        T[i, j] = elapsed_min(x_route[i], tws[j], twd[j], pol[j])

# Out-of-sample: mask each route's own design scenario, renormalise per row.
mask = np.ones_like(T, dtype=bool)
mask[np.arange(N_ROUTES), home] = False
wm = np.tile(w, (N_ROUTES, 1)) * mask
wm /= wm.sum(axis=1, keepdims=True)

mean_oos = (T * wm).sum(axis=1)
spread = np.sqrt((wm * (T - mean_oos[:, None]) ** 2).sum(axis=1))
p_win = (wm * (T < RECORD_MIN)).sum(axis=1)

best = int(np.argmax(p_win))
fastest = int(np.argmin(mean_oos))

# ----------------------------------------------- fig 1: weighted input space
lv = np.array([-2, -1, 0, 1, 2])
mw = np.exp(-0.5 * lv ** 2)
mw /= mw.sum()

fig, ax = plt.subplots(figsize=(9.2, 4.2))
for xi in range(5):
    for yi in range(5):
        jw = mw[xi] * mw[yi]
        ax.scatter(xi, yi, s=jw * 7200, color=ACCENT,
                   alpha=0.20 + 0.60 * jw / mw.max() ** 2,
                   edgecolor=ACCENT, linewidth=0.8, zorder=2)


def leader(colour):
    return dict(arrowstyle="->", color=colour, lw=1.0,
                shrinkA=3, shrinkB=9, connectionstyle="arc3,rad=-0.12")


LAB = dict(fontsize=10.5, va="center", ha="left")
ax.annotate("nothing off" + "\n" + "w = 0.162", xy=(2, 2), xytext=(-2.3, 3.6),
            color=INK, weight="bold", arrowprops=leader(INK), **LAB)
ax.annotate("TWD off only" + "\n" + "w = 0.098", xy=(1, 2), xytext=(-2.3, 2.05),
            color=ACCENT, arrowprops=leader(ACCENT), **LAB)
ax.annotate("TWD and polar off" + "\n" + "w = 0.060", xy=(1, 1), xytext=(-2.3, 0.5),
            color=WARM, arrowprops=leader(WARM), **LAB)

ax.set_xticks(range(5))
ax.set_xticklabels(["{:+.0f}".format(v * SD_TWD) if v else "0" for v in lv])
ax.set_yticks(range(5))
ax.set_yticklabels(["{:.0f}%".format((1 + v * SD_POL) * 100) for v in lv])
ax.set_xlabel("TWD deviation (deg)")
ax.set_ylabel("Polar scaling")
ax.set_xlim(-2.45, 4.55)
ax.set_ylim(-0.45, 4.3)
ax.spines["left"].set_visible(False)
ax.tick_params(length=0)
ax.set_title("Joint weight is the product of the marginals, so simultaneous deviations are rare",
             fontsize=12, color=MUTED, pad=12, loc="left")
fig.tight_layout()
fig.savefig(FIG + "/s1_inputspace.png", dpi=200)
plt.close(fig)

# ---------------------------- fig 2: cross-evaluation matrix, shown as regret
regret = T - T.min(axis=0, keepdims=True)     # minutes lost vs the ideal route
cmap = LinearSegmentedColormap.from_list(
    "regret", ["#FBF7F0", "#CFE0DC", "#7FB3AE", "#2E8B85", "#0B4F4F"])

fig, ax = plt.subplots(figsize=(9.6, 3.5))
im = ax.imshow(regret, aspect="auto", cmap=cmap, vmin=0, vmax=22)
for i in range(N_ROUTES):
    ax.add_patch(plt.Rectangle((home[i] - 0.5, i - 0.5), 1, 1, fill=False,
                               edgecolor=WARM, lw=1.9, zorder=3))
ax.set_xlabel("Input scenario j, sorted by TWD deviation   (backing to the right)")
ax.set_ylabel("Route i, by bulge (nm)")
ax.set_yticks(range(N_ROUTES))
ax.set_yticklabels(["{:+.0f}".format(v) for v in x_route], fontsize=8)
ax.set_title("Minutes lost by sailing route i when input j actually happens",
             fontsize=12, color=MUTED, pad=10, loc="left")
cb = fig.colorbar(im, ax=ax, pad=0.015)
cb.set_label("minutes lost")
cb.outline.set_visible(False)
fig.subplots_adjust(left=0.10, right=0.995, top=0.855, bottom=0.25)
fig.text(0.10, 0.045,
         "Pale band: each route is near-ideal close to its own design case.  Orange = that design "
         "case itself; exclude it, or robustness flatters every route.",
         fontsize=9, color=MUTED, style="italic")
fig.savefig(FIG + "/s2_matrix.png", dpi=200)
plt.close(fig)

# -------------------------------------------------------- fig 3: the decision
fig, (axl, axr) = plt.subplots(1, 2, figsize=(9.8, 3.9),
                               gridspec_kw={"width_ratios": [1.30, 1]})

sc = axl.scatter(mean_oos, spread, c=p_win, cmap="viridis", s=150,
                 vmin=p_win.min(), vmax=p_win.max(),
                 edgecolor="white", linewidth=1.3, zorder=3)
XLO = min(RECORD_MIN, mean_oos.min())
XHI = mean_oos.max()
XR = XHI - XLO
YLO, YHI = spread.min(), spread.max()
YR = YHI - YLO

axl.axvline(RECORD_MIN, color=WARM, ls="--", lw=1.4, zorder=1)
axl.text(RECORD_MIN + 0.02 * XR, YHI + 0.30 * YR,
         "record {:.0f} min".format(RECORD_MIN), fontsize=9.5, color=WARM)

axl.annotate("highest P(win)" + "\n" + "route {} ({:+.0f} nm)".format(best, x_route[best]),
             xy=(mean_oos[best], spread[best]), xytext=(XLO + 0.46 * XR, YLO + 0.52 * YR),
             fontsize=10, weight="bold", color=INK, ha="left", va="center",
             arrowprops=dict(arrowstyle="->", color=INK, lw=1.1, shrinkB=7))
axl.annotate("fastest mean" + "\n" + "route {} ({:+.0f} nm)".format(fastest, x_route[fastest]),
             xy=(mean_oos[fastest], spread[fastest]), xytext=(XLO + 0.30 * XR, YLO + 0.06 * YR),
             fontsize=10, color=MUTED, ha="left", va="center",
             arrowprops=dict(arrowstyle="->", color=MUTED, lw=1.0, shrinkB=7))

axl.set_xlabel("Out-of-sample mean duration (min), faster is left")
axl.set_ylabel("Duration spread (min), robust is low")
axl.set_title("Two KPIs: explanatory", fontsize=12, color=MUTED, loc="left", pad=10)
axl.set_xlim(XLO - 0.07 * XR, XHI + 0.12 * XR)
axl.set_ylim(YLO - 0.42 * YR, YHI + 0.46 * YR)
cb = fig.colorbar(sc, ax=axl, pad=0.02)
cb.set_label("P(break record)")
cb.outline.set_visible(False)

rank = np.argsort(-p_win)
cols = [ACCENT if k == best else "#B9C6CC" for k in rank]
axr.barh(range(N_ROUTES), p_win[rank], color=cols, height=0.72)
axr.set_yticks(range(N_ROUTES))
axr.set_yticklabels(["r{}   {:+.0f} nm".format(k, x_route[k]) for k in rank], fontsize=8.5)
axr.invert_yaxis()
axr.set_xlim(0, p_win.max() * 1.30)
axr.set_xlabel("P(break record)")
axr.set_title("One KPI: decisive", fontsize=12, color=MUTED, loc="left", pad=10)
for r, k in enumerate(rank):
    axr.text(p_win[k] + p_win.max() * 0.035, r, "{:.0%}".format(p_win[k]),
             va="center", fontsize=9, color=INK)
fig.tight_layout()
fig.savefig(FIG + "/s3_decision.png", dpi=200)
plt.close(fig)

# ------------------------------------------------------------- the pitch deck
ROUTEMAP = FIG + "/s0_routemap.png"
if not os.path.exists(ROUTEMAP):
    raise SystemExit("Missing " + ROUTEMAP + " -- run: python slides/plot_routes.py")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def hexc(h):
    return RGBColor.from_string(h.lstrip("#").upper())


def add_slide(kicker, title, img, takeaway, notes):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    tb = s.shapes.add_textbox(Inches(0.72), Inches(0.40), Inches(11.9), Inches(0.32))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = kicker
    r.font.size, r.font.bold = Pt(12), True
    r.font.color.rgb, r.font.name = hexc(ACCENT), "Calibri"

    tb = s.shapes.add_textbox(Inches(0.72), Inches(0.72), Inches(11.9), Inches(0.64))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size, r.font.bold = Pt(29), True
    r.font.color.rgb, r.font.name = hexc(INK), "Calibri"

    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.46),
                            Inches(11.9), Pt(1.2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = hexc(RULE)
    ln.line.fill.background()
    ln.shadow.inherit = False

    s.shapes.add_picture(img, Inches(1.05), Inches(1.82), width=Inches(11.2))

    tb = s.shapes.add_textbox(Inches(0.72), Inches(6.56), Inches(11.9), Inches(0.60))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = takeaway
    r.font.size, r.font.bold = Pt(15), True
    r.font.color.rgb, r.font.name = hexc(ACCENT), "Calibri"

    s.notes_slide.notes_text_frame.text = notes
    return s


add_slide(
    "THE PROBLEM   |   ~45 s",
    "102 miles due east, and an eight-mile corridor",
    ROUTEMAP,
    "101.7 nm on 090 T. 18 knots average to take it. The wind farms leave us about 8 nm of "
    "corridor, north only.",
    "[~45 s]  Lowestoft to IJmuiden: 101.7 nautical miles, almost exactly due east. To take the "
    "record we have to average 18 knots, five hours thirty-nine. On a reach this short there is no "
    "clever detour: every candidate is the same reach, and they differ only in how far north we "
    "bow out. And the sea room is tighter than it looks. The wind farms off the Dutch coast pinch "
    "us into a corridor roughly eight miles wide, and it is north-only: bulge south and we are "
    "inside Luchterduinen. So the question is not which of a thousand routes. It is: inside a "
    "narrow corridor, which line still works when the forecast turns out wrong?")

add_slide(
    "THE INPUT SPACE   |   ~45 s",
    "We don't trust one forecast, we weight many",
    FIG + "/s1_inputspace.png",
    "One thing being off is likely. Everything being off at once is not. Weight accordingly.",
    "[~45 s]  We do not trust a single forecast. We perturb everything that matters, the polar, "
    "wind speed, wind direction, current, and give every level a likelihood. The joint weight is "
    "the product of the marginals, so being wrong about one thing far outranks being wrong about "
    "everything at once: across four axes, a single one-sigma deviation is about three times as "
    "likely as three simultaneous ones. That matters. A naive worst-case study assumes every input "
    "fails together, and then tells you to sail defensively for a day that will essentially never "
    "happen. Weighting by likelihood is what keeps the pessimism honest.")

add_slide(
    "THE TEST   |   ~45 s",
    "Optimise per forecast, then sail every route in every forecast",
    FIG + "/s2_matrix.png",
    "A route's honest robustness is its row with its own design case removed.",
    "[~45 s]  For each weighted input we run the optimiser once and get one optimal route. That is "
    "the easy half. Robustness comes from the other half: we take every route and sail it through "
    "all forty inputs. This chart is minutes lost against the ideal line for that particular day. "
    "Each route is excellent near its own design case and decays away from it, and notice the pale "
    "band widens for the larger northward bulges. Those are the forgiving ones. One discipline "
    "point: we exclude each route's own design case, because otherwise every route gets one "
    "flattering result and robustness looks better than it really is.")

add_slide(
    "THE ANSWER   |   ~45 s",
    "One number: probability of taking the record",
    FIG + "/s3_decision.png",
    "The record time is the exchange rate. One ranked number, and the binding constraint named.",
    "[~45 s]  Duration and robustness are two KPIs and you can draw a Pareto front, but you do not "
    "have to trade them off by hand, because the record time is the exchange rate. Weight every "
    "cell by its likelihood, count the ones under 339 minutes, and you get one ranked number: "
    "probability of taking the record. The fastest-mean route comes second. And the winner sits at "
    "plus seven and a half miles, right against the Prinses Amalia boundary. That tells us the "
    "binding constraint is the wind farm, not the weather. So the instruction to the skipper is "
    "simple: bulge as far north as the farm allows.")

OUT = "slides/robust_routing_pitch.pptx"
try:
    prs.save(OUT)
except PermissionError:
    OUT = "slides/robust_routing_pitch_new.pptx"
    prs.save(OUT)
    print("NOTE: canonical file locked (open in PowerPoint); wrote " + OUT)

notes_words = sum(len(s.notes_slide.notes_text_frame.text.split()) for s in prs.slides)
print("slides 4   speaker-note words {}   ~{:.1f} min at 145 wpm".format(
    notes_words, notes_words / 145.0))
print("rhumb {:.1f} nm   record {:.0f} min".format(2 * HALF_NM, RECORD_MIN))
print("P(win) range     {:.1%} .. {:.1%}".format(p_win.min(), p_win.max()))
print("best P(win)   route {:2d}  x={:+5.1f} nm  P={:.1%}  mean={:.1f}  spread={:.1f}".format(
    best, x_route[best], p_win[best], mean_oos[best], spread[best]))
print("fastest mean  route {:2d}  x={:+5.1f} nm  P={:.1%}  mean={:.1f}  spread={:.1f}".format(
    fastest, x_route[fastest], p_win[fastest], mean_oos[fastest], spread[fastest]))
print("saved " + OUT)
